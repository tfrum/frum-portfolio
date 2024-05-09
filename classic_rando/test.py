import sys
print(sys.version)


# mergeTwoLists = lambda l1, l2, dummy=ListNode(): (lambda f, dummy: (f(f, dummy, l1, l2), dummy.next)[1])(lambda rec, node, l1, l2: node if not l1 and not l2 else (rec(rec, node.__setattr__('next', l1 if not l2 or (l1 and l1.val < l2.val) else l2) or node.next, l1.next if not l2 or (l1 and l1.val < l2.val) else l1, l2 if not l2 or (l1 and l1.val < l2.val) else l2.next)), dummy)
from pptx import Presentation
from pptx.util import Inches

def create_presentation():
    # Create a presentation object
    prs = Presentation()

    # Add a slide with the title and content layout
    slide_layout = prs.slide_layouts[1]  # 0 is title slide, 1 is title and content
    slide = prs.slides.add_slide(slide_layout)

    # Set title and subtitle
    title = slide.shapes.title
    title.text = "Digital Twins"

    content = slide.placeholders[1]
    content.text = "Introduction to Digital Twins"

    # Add bullet points
    text_frame = content.text_frame
    p = text_frame.add_paragraph()
    p.text = "A digital twin is a virtual model designed to accurately reflect a physical object."
    p = text_frame.add_paragraph()
    p.text = "Used for analysis, to improve operational efficiency, or to predict future performance."
    p = text_frame.add_paragraph()
    p.text = "Common in industries like manufacturing, automotive, and healthcare."

    # Add an image (ensure the image file exists)
    img_path = 'path_to_image.jpg'  # Specify the path to your image file
    top = Inches(2)
    left = Inches(5)  # Distance from left of slide
    height = Inches(3)  # Height of image
    #pic = slide.shapes.add_picture(img_path, left, top, height=height)

    # Save the presentation
    prs.save('Digital_Twins_Presentation.pptx')

if __name__ == "__main__":
    create_presentation()
    print("Presentation created successfully.")
